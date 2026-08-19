"""Document text extraction.

Every format is normalised into the *same* list of text blocks, which then go
through the project's existing chunking (1000 / 100 overlap) and the existing
768-dimensional embedding model. Neither the chunker nor the embedder changed.

    PDF / DOC / DOCX / PPT / PPTX / TXT / MD
                    |
        format-specific extraction (+ OCR for scanned PDFs)
                    |
              normalised blocks   [{text, metadata}]
                    |
        RecursiveCharacterTextSplitter (1000 / 100)
                    |
                  chunks -> embeddings -> pgvector

Metadata is preserved per format so citations can point somewhere real:
PDF -> page, DOCX -> section/paragraph, PPTX -> slide, TXT/MD -> line/section.
"""

from __future__ import annotations

import os
import subprocess
import tempfile
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter

import capabilities
import settings

# Extensions the code knows how to handle; whether they *work* here is
# answered by capabilities.supported_extensions().
KNOWN_EXTENSIONS = {".pdf", ".txt", ".md", ".docx", ".doc", ".pptx", ".ppt"}

# A PDF page with less than this many characters of embedded text is treated as
# scanned and sent to OCR. Deliberately low so normal text PDFs are never OCRed.
OCR_CHAR_THRESHOLD = 40
OCR_DPI = 200

# Pages read before the PDF reader is reopened to release its object cache.
# Peak memory scales with this, not with the length of the document.
PDF_PAGE_WINDOW = max(1, int(os.getenv("PDF_PAGE_WINDOW", "50")))

# How often a text page reports progress. OCR pages always report, being slow.
PROGRESS_PAGE_INTERVAL = max(1, int(os.getenv("PROGRESS_PAGE_INTERVAL", "10")))


def _report(progress, **event) -> None:
    """Send one progress event, if anybody is listening.

    Every field is a real count taken from the work actually done; nothing here
    is estimated, and a caller that cannot measure something simply omits it.
    """
    if progress is None:
        return
    try:
        progress(event)
    except Exception:  # noqa: BLE001
        # Progress reporting must never break an ingestion.
        pass


class ExtractionError(Exception):
    """Raised with a user-facing message when a document cannot be read."""

    def __init__(self, message: str, code: str = "EXTRACTION_FAILED"):
        super().__init__(message)
        self.code = code


class UnsupportedFormat(ExtractionError):
    def __init__(self, message: str):
        super().__init__(message, "UNSUPPORTED_FILE_TYPE")


# --------------------------------------------------------------------------
# Legacy Office conversion (.doc / .ppt) via LibreOffice
# --------------------------------------------------------------------------

def _convert_with_soffice(path: Path, target_ext: str) -> Path:
    """Convert a legacy binary Office file to its OOXML equivalent.

    LibreOffice is a real system dependency; if it is absent the caller gets a
    clear message rather than an empty document.
    """
    soffice = capabilities.libreoffice_path()
    if not soffice:
        raise UnsupportedFormat(
            f"{path.suffix.upper()} files need LibreOffice on the server to be read. "
            f"Install LibreOffice and set LIBREOFFICE_PATH, or save the file as "
            f"{target_ext.upper()} and upload it again."
        )

    outdir = Path(tempfile.mkdtemp(prefix="soffice-"))
    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--norestore",
                # Macros are never executed: the converter runs headless and
                # LibreOffice's default macro security refuses them.
                "--convert-to",
                target_ext.lstrip("."),
                "--outdir",
                str(outdir),
                str(path),
            ],
            capture_output=True,
            timeout=180,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise ExtractionError("Converting this document timed out.", "CONVERSION_TIMEOUT") from exc

    converted = list(outdir.glob(f"*{target_ext}"))
    if not converted:
        detail = (result.stderr or b"").decode("utf-8", "replace")[:200]
        raise ExtractionError(
            "LibreOffice could not convert this document. It may be corrupted or password protected."
            + (f" ({detail})" if detail else ""),
            "CONVERSION_FAILED",
        )

    return converted[0]


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------

def _ocr_pdf_pages(path: Path, page_numbers: list[int], progress=None) -> dict[int, str]:
    """OCR specific pages. Rendering uses PyMuPDF, so poppler is not needed."""
    if not page_numbers:
        return {}

    if not capabilities.ocr_available():
        raise ExtractionError(
            "This PDF appears to be scanned and needs OCR, but the Tesseract OCR engine is not "
            "installed on the server. Install Tesseract, or upload a PDF that has selectable text.",
            "OCR_UNAVAILABLE",
        )

    import pymupdf
    import pytesseract
    from PIL import Image

    pytesseract.pytesseract.tesseract_cmd = capabilities.tesseract_path()

    results: dict[int, str] = {}
    with pymupdf.open(str(path)) as document:
        for position, page_index in enumerate(page_numbers):
            _report(
                progress,
                stage="ocr",
                page=page_index + 1,
                pages=len(page_numbers),
                ocrPages=position + 1,
            )
            page = document[page_index]
            pixmap = page.get_pixmap(dpi=OCR_DPI)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            try:
                results[page_index] = pytesseract.image_to_string(image).strip()
            except Exception as exc:  # noqa: BLE001
                raise ExtractionError(f"OCR failed on page {page_index + 1}: {exc}", "OCR_FAILED")

    return results


def _ocr_one_page(document, page_index: int) -> str:
    """OCR a single already-open PDF page. Nothing is retained afterwards."""
    import pytesseract
    from PIL import Image

    page = document[page_index]
    pixmap = page.get_pixmap(dpi=OCR_DPI)
    try:
        image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
        try:
            return pytesseract.image_to_string(image).strip()
        finally:
            image.close()
    finally:
        # Release the rendered bitmap immediately; a 200 DPI page is ~10 MB.
        del pixmap


def _open_pdf(path: Path):
    """Opens a PDF for reading, unlocking it if it has an empty user password."""
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                raise ExtractionError(
                    "This PDF is password protected. Remove the password and upload it again.",
                    "ENCRYPTED",
                )
        return reader
    except ExtractionError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            "Could not read this PDF. It may be corrupted or not a valid PDF file.",
            "PARSE_FAILED",
        ) from exc


def _close_pdf(reader) -> None:
    """Releases the file handle and every cached object held by the reader."""
    try:
        reader.close()
    except Exception:  # noqa: BLE001
        # Older pypdf has no close(); drop the caches by hand instead.
        for attribute in ("resolved_objects", "xref", "xref_objStm"):
            try:
                getattr(reader, attribute).clear()
            except Exception:  # noqa: BLE001
                pass


def _iter_pdf_blocks(path: Path, name: str, progress=None, info: dict | None = None):
    """Yield one block per page, reading the PDF a page at a time.

    Peak memory is one window of pages rather than the whole document, so a
    500-page file costs about the same as a 50-page one. Pages whose embedded
    text is too short are OCRed individually, preserving the "never OCR a page
    that already has text" behaviour.
    """
    info = info if info is not None else {}

    reader = _open_pdf(path)
    try:
        page_count = len(reader.pages)
    except Exception as exc:  # noqa: BLE001
        _close_pdf(reader)
        raise ExtractionError(
            "Could not read this PDF. It may be corrupted or not a valid PDF file.",
            "PARSE_FAILED",
        ) from exc

    if not page_count:
        _close_pdf(reader)
        raise ExtractionError("This PDF has no pages.", "EMPTY_DOCUMENT")

    info["pageCount"] = page_count
    ocr_pages = 0
    ocr_document = None

    # The page count is known before any page is read, so the UI can show real
    # "page N of M" progress for the whole extraction rather than a spinner.
    _report(progress, stage="extracting", page=0, pages=page_count)

    try:
        for index in range(page_count):
            # pypdf caches every object it resolves for the life of the reader,
            # so memory grows with the number of pages read no matter how the
            # output is batched. Reopening at each window boundary drops that
            # cache; nothing already emitted is needed again.
            if index and index % PDF_PAGE_WINDOW == 0:
                _close_pdf(reader)
                reader = _open_pdf(path)

            try:
                text = (reader.pages[index].extract_text() or "").strip()
            except Exception:
                text = ""

            # Scanned page: OCR just this one.
            if len(text) < OCR_CHAR_THRESHOLD:
                if not capabilities.ocr_available():
                    raise ExtractionError(
                        "This PDF appears to be scanned and needs OCR, but the Tesseract OCR "
                        "engine is not installed on the server. Install Tesseract, or upload a "
                        "PDF that has selectable text.",
                        "OCR_UNAVAILABLE",
                    )

                import pymupdf
                import pytesseract

                if ocr_document is None:
                    pytesseract.pytesseract.tesseract_cmd = capabilities.tesseract_path()
                    ocr_document = pymupdf.open(str(path))

                # OCR is the slow path, so every page is reported: a 180-page
                # scan otherwise looks frozen for several minutes.
                _report(
                    progress,
                    stage="ocr",
                    page=index + 1,
                    pages=page_count,
                    ocrPages=ocr_pages + 1,
                )

                try:
                    recognised = _ocr_one_page(ocr_document, index)
                except ExtractionError:
                    raise
                except Exception as exc:  # noqa: BLE001
                    raise ExtractionError(
                        f"OCR failed on page {index + 1}: {exc}", "OCR_FAILED"
                    ) from exc

                if recognised:
                    text = recognised
                    ocr_pages += 1

            if text.strip():
                yield {"text": text, "metadata": {"page": index + 1}}

            # Text pages are fast, so reporting each one would be noise. One
            # tick per window is enough to keep the page counter moving.
            if (index + 1) % PROGRESS_PAGE_INTERVAL == 0 or index + 1 == page_count:
                _report(progress, stage="extracting", page=index + 1, pages=page_count)
    finally:
        _close_pdf(reader)
        if ocr_document is not None:
            ocr_document.close()

    info["usedOcr"] = ocr_pages > 0
    info["ocrPages"] = ocr_pages


def _extract_pdf(path: Path, name: str, progress=None) -> tuple[list[dict], dict]:
    """Eager wrapper kept for callers that want the whole document at once."""
    info: dict = {}
    blocks = list(_iter_pdf_blocks(path, name, progress, info))
    info.setdefault("usedOcr", False)
    info.setdefault("ocrPages", 0)
    return blocks, info


# --------------------------------------------------------------------------
# Word
# --------------------------------------------------------------------------

def _extract_docx(path: Path, name: str) -> tuple[list[dict], dict]:
    """Paragraphs, headings, lists and tables - never treated as an image."""
    try:
        import docx
    except ImportError as exc:
        raise UnsupportedFormat(
            "DOCX support requires python-docx on the server (pip install python-docx)."
        ) from exc

    try:
        document = docx.Document(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            "Could not read this Word document. It may be corrupted, password protected, "
            "or saved in an unsupported format.",
            "PARSE_FAILED",
        ) from exc

    blocks: list[dict] = []
    section = None
    paragraph_number = 0

    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if not text:
            continue

        paragraph_number += 1
        style = (paragraph.style.name if paragraph.style else "") or ""

        # Track the current heading so chunks can say which section they are from.
        if style.startswith("Heading") or style == "Title":
            section = text
            blocks.append(
                {
                    "text": f"# {text}" if style == "Title" else f"## {text}",
                    "metadata": {"paragraph": paragraph_number, "section": section, "heading": True},
                }
            )
            continue

        # Keep list structure readable in the extracted text.
        prefix = "- " if "List" in style else ""
        blocks.append(
            {
                "text": f"{prefix}{text}",
                "metadata": {"paragraph": paragraph_number, "section": section},
            }
        )

    # Tables carry real content in a lot of documents.
    for table_index, table in enumerate(document.tables, 1):
        rows = []
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append(
                {
                    "text": "\n".join(rows),
                    "metadata": {"table": table_index, "section": section},
                }
            )

    return blocks, {"paragraphs": paragraph_number, "tables": len(document.tables)}


def _extract_doc(path: Path, name: str) -> tuple[list[dict], dict]:
    converted = _convert_with_soffice(path, ".docx")
    blocks, meta = _extract_docx(converted, name)
    return blocks, {**meta, "convertedFrom": ".doc"}


# --------------------------------------------------------------------------
# PowerPoint
# --------------------------------------------------------------------------

def _extract_pptx(path: Path, name: str) -> tuple[list[dict], dict]:
    """Titles, text boxes, bullets, tables and speaker notes, per slide."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise UnsupportedFormat(
            "PPTX support requires python-pptx on the server (pip install python-pptx)."
        ) from exc

    try:
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(
            "Could not read this PowerPoint file. It may be corrupted or password protected.",
            "PARSE_FAILED",
        ) from exc

    blocks: list[dict] = []
    slide_count = 0

    for slide_number, slide in enumerate(presentation.slides, 1):
        slide_count += 1
        title = None
        parts: list[str] = []

        for shape in slide.shapes:
            # Titles first so the slide has a heading.
            if shape == slide.shapes.title and shape.has_text_frame:
                title = (shape.text_frame.text or "").strip()
                continue

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if not text:
                        continue
                    # Indent level makes the bullet hierarchy readable.
                    indent = "  " * max(getattr(paragraph, "level", 0) or 0, 0)
                    parts.append(f"{indent}- {text}")

            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        body = "\n".join(parts).strip()
        if not (title or body or notes):
            continue

        segments = []
        if title:
            segments.append(f"# {title}")
        if body:
            segments.append(body)
        if notes:
            segments.append(f"Speaker notes: {notes}")

        blocks.append(
            {
                "text": "\n".join(segments),
                "metadata": {
                    "slide": slide_number,
                    **({"title": title} if title else {}),
                    **({"hasNotes": True} if notes else {}),
                },
            }
        )

    return blocks, {"slideCount": slide_count}


def _extract_ppt(path: Path, name: str) -> tuple[list[dict], dict]:
    converted = _convert_with_soffice(path, ".pptx")
    blocks, meta = _extract_pptx(converted, name)
    return blocks, {**meta, "convertedFrom": ".ppt"}


# --------------------------------------------------------------------------
# Plain text and Markdown
# --------------------------------------------------------------------------

def _read_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ExtractionError("Could not decode this text file.", "DECODE_FAILED")


def _extract_txt(path: Path, name: str) -> tuple[list[dict], dict]:
    """Split on blank lines so paragraph structure survives; track line numbers."""
    text = _read_text(path)
    blocks: list[dict] = []

    line_number = 1
    for paragraph in text.split("\n\n"):
        lines = paragraph.split("\n")
        cleaned = paragraph.strip()
        if cleaned:
            blocks.append({"text": cleaned, "metadata": {"line": line_number}})
        line_number += len(lines) + 1

    return blocks, {"lines": text.count("\n") + 1}


def _extract_markdown(path: Path, name: str) -> tuple[list[dict], dict]:
    """Keep Markdown structure and record the heading each block sits under."""
    text = _read_text(path)
    blocks: list[dict] = []

    section: str | None = None
    buffer: list[str] = []
    buffer_line = 1
    in_code_fence = False
    headings = 0

    def flush(start_line: int) -> None:
        content = "\n".join(buffer).strip()
        if content:
            blocks.append(
                {
                    "text": content,
                    "metadata": {"line": start_line, **({"section": section} if section else {})},
                }
            )
        buffer.clear()

    for index, line in enumerate(text.split("\n"), 1):
        if line.lstrip().startswith("```"):
            in_code_fence = not in_code_fence
            buffer.append(line)
            continue

        # A heading outside a code fence starts a new section.
        if not in_code_fence and line.lstrip().startswith("#"):
            flush(buffer_line)
            section = line.lstrip("#").strip() or section
            headings += 1
            buffer_line = index
            buffer.append(line)
            continue

        if not buffer:
            buffer_line = index
        buffer.append(line)

    flush(buffer_line)
    return blocks, {"headings": headings, "lines": text.count("\n") + 1}


# --------------------------------------------------------------------------
# Dispatch
# --------------------------------------------------------------------------

EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_doc,
    ".pptx": _extract_pptx,
    ".ppt": _extract_ppt,
    ".txt": _extract_txt,
    ".md": _extract_markdown,
}


def extract_blocks(file_path: str, display_name: str | None = None, progress=None):
    """Format-specific extraction into normalised blocks."""
    path = Path(file_path)
    if not path.exists():
        raise ExtractionError("The uploaded file could not be found on the server.", "FILE_MISSING")

    extension = path.suffix.lower()
    if extension not in KNOWN_EXTENSIONS:
        raise UnsupportedFormat(
            f"Unsupported file type '{extension or 'unknown'}'. "
            f"Supported: {', '.join(sorted(KNOWN_EXTENSIONS))}."
        )

    if path.stat().st_size == 0:
        raise ExtractionError("This file is empty.", "EMPTY_DOCUMENT")

    name = display_name or path.name
    extractor = EXTRACTORS[extension]

    if extension == ".pdf":
        return extractor(path, name, progress)
    return extractor(path, name)


def iter_blocks(file_path: str, display_name: str | None = None, progress=None, info=None):
    """Yield normalised blocks lazily.

    PDFs stream a page at a time, which is where the memory matters. The other
    formats are parsed by libraries that load the whole file anyway, so their
    blocks are produced up front and then yielded one by one - still avoiding a
    second full copy downstream.
    """
    path, name, extension = _prepare(file_path, display_name)
    info = info if info is not None else {}

    if extension == ".pdf":
        yield from _iter_pdf_blocks(path, name, progress, info)
        return

    blocks, block_info = EXTRACTORS[extension](path, name)
    info.update(block_info)

    # These formats are parsed in one go, so the totals are known immediately.
    # Reporting them lets the UI say how much there is even though extraction
    # itself was too fast to show progress for.
    _report(
        progress,
        stage="extracting",
        blocks=len(blocks),
        **({"pages": block_info["slideCount"]} if "slideCount" in block_info else {}),
    )

    for index, block in enumerate(blocks, 1):
        yield block
        _report(progress, stage="extracting", block=index, blocks=len(blocks))


def iter_chunk_batches(
    file_path: str,
    display_name: str | None = None,
    batch_size: int = 64,
    progress=None,
    info=None,
):
    """Yield lists of at most `batch_size` chunks, ready to embed.

    This is the memory-bounded entry point used by ingestion: at no point does
    more than one batch of chunks exist, so a 1000-page document costs roughly
    the same as a one-page one. Chunking itself is unchanged
    (RecursiveCharacterTextSplitter at the configured 1000 / 100).
    """
    path = Path(file_path)
    name = display_name or path.name
    info = info if info is not None else {}

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
    )

    batch: list[dict] = []
    produced = 0

    for block in iter_blocks(file_path, name, progress, info):
        for piece in splitter.split_text(block["text"]):
            content = piece.strip()
            if not content:
                continue

            batch.append(
                {
                    "content": content,
                    "metadata": {
                        "document_name": name,
                        "chunk_index": produced,
                        **block["metadata"],
                    },
                }
            )
            produced += 1

            if len(batch) >= batch_size:
                yield batch
                # Hand ownership to the consumer and start a fresh list, so the
                # batch just yielded can be collected as soon as it is used.
                batch = []

    if batch:
        yield batch

    info["blocks"] = info.get("blocks", 0)
    info["chunks"] = produced

    if produced == 0:
        raise ExtractionError(
            "Could not extract readable text from this document. It may contain only images, "
            "or its text may not be selectable.",
            "NO_TEXT_EXTRACTED",
        )


def _prepare(file_path: str, display_name: str | None):
    """Shared validation: existence, known type, non-empty."""
    path = Path(file_path)
    if not path.exists():
        raise ExtractionError("The uploaded file could not be found on the server.", "FILE_MISSING")

    extension = path.suffix.lower()
    if extension not in KNOWN_EXTENSIONS:
        raise UnsupportedFormat(
            f"Unsupported file type '{extension or 'unknown'}'. "
            f"Supported: {', '.join(sorted(KNOWN_EXTENSIONS))}."
        )

    if path.stat().st_size == 0:
        raise ExtractionError("This file is empty.", "EMPTY_DOCUMENT")

    return path, display_name or path.name, extension


def extract_chunks(file_path: str, display_name: str | None = None, progress=None) -> dict:
    """Full extraction + chunking, materialised.

    Thin wrapper over iter_chunk_batches for callers that genuinely want every
    chunk at once (the tests, the legacy CLI). Ingestion uses the streaming
    generator instead so it never holds the whole document.
    """
    info: dict = {}
    chunks: list[dict] = []
    for batch in iter_chunk_batches(file_path, display_name, 64, progress, info):
        chunks.extend(batch)

    return {"chunks": chunks, "info": {**info, "blocks": info.get("blocks", len(chunks))}}
