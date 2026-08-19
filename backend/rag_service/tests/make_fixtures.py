"""Generates real test documents for every supported (and unsupported) format.

Everything is produced programmatically so the fixtures are reproducible and no
binary blobs need to live in the repository.
"""

from __future__ import annotations

import sys
from pathlib import Path

FIXTURES = Path(__file__).resolve().parent / "fixtures"


def make_txt() -> Path:
    path = FIXTURES / "sample.txt"
    path.write_text(
        "Quarterly Report\n"
        "\n"
        "Revenue grew by 14 percent this quarter, driven mainly by the retail\n"
        "segment and a recovery in European orders.\n"
        "\n"
        "Operating costs stayed flat despite inflationary pressure on logistics.\n"
        "\n"
        "The board approved a dividend of 0.42 per share.\n",
        encoding="utf-8",
    )
    return path


def make_md() -> Path:
    path = FIXTURES / "sample.md"
    path.write_text(
        "# Vector Databases\n"
        "\n"
        "A vector database stores embeddings and supports similarity search.\n"
        "\n"
        "## Indexing strategies\n"
        "\n"
        "- HNSW builds a navigable small-world graph\n"
        "- IVF partitions the space into cells\n"
        "\n"
        "## Example query\n"
        "\n"
        "```sql\n"
        "SELECT id FROM chunks ORDER BY embedding <=> $1 LIMIT 5;\n"
        "```\n"
        "\n"
        "Cosine distance is appropriate for normalised vectors.\n",
        encoding="utf-8",
    )
    return path


def make_pdf_text() -> Path:
    """A normal PDF with a real, selectable text layer."""
    import pymupdf

    path = FIXTURES / "text.pdf"
    document = pymupdf.open()
    for number, body in enumerate(
        [
            "Gated Recurrent Units\nA GRU uses an update gate and a reset gate to control "
            "how much past state is carried forward.",
            "Training\nGRUs are trained with backpropagation through time and are cheaper "
            "than LSTM cells because they have fewer parameters.",
        ],
        1,
    ):
        page = document.new_page()
        page.insert_text((72, 96), f"Page {number}", fontsize=16)
        page.insert_textbox(pymupdf.Rect(72, 120, 520, 400), body, fontsize=11)
    document.save(str(path))
    document.close()
    return path


def _text_image(text: str, size=(1240, 400)):
    """Renders text to a bitmap - no text layer, so only OCR can read it."""
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("arial.ttf", 42)
    except OSError:
        font = ImageFont.load_default(size=42)
    draw.multiline_text((40, 40), text, fill="black", font=font, spacing=18)
    return image


def make_pdf_scanned() -> Path:
    """A 'scanned' PDF: page content is an image, with no extractable text."""
    import io

    import pymupdf

    path = FIXTURES / "scanned.pdf"
    image = _text_image("SCANNED INVOICE 4471\nTotal due 1250 USD\nPayable within 30 days")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")

    document = pymupdf.open()
    page = document.new_page(width=620, height=200)
    page.insert_image(pymupdf.Rect(0, 0, 620, 200), stream=buffer.getvalue())
    document.save(str(path))
    document.close()
    return path


def make_pdf_mixed() -> Path:
    """Page 1 has real text, page 2 is an image only."""
    import io

    import pymupdf

    path = FIXTURES / "mixed.pdf"
    document = pymupdf.open()

    page = document.new_page()
    page.insert_textbox(
        pymupdf.Rect(72, 96, 520, 400),
        "Chapter One\nThis page contains ordinary selectable text that pypdf can read "
        "directly without any optical character recognition.",
        fontsize=12,
    )

    image = _text_image("APPENDIX SCAN\nFigure 2 shows the measured throughput")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    page2 = document.new_page(width=620, height=200)
    page2.insert_image(pymupdf.Rect(0, 0, 620, 200), stream=buffer.getvalue())

    document.save(str(path))
    document.close()
    return path


def make_docx() -> Path:
    import docx

    path = FIXTURES / "sample.docx"
    document = docx.Document()
    document.add_heading("Machine Learning Overview", level=1)
    document.add_paragraph(
        "Supervised learning fits a function from labelled examples to a target variable."
    )
    document.add_heading("Common algorithms", level=2)
    document.add_paragraph("Gradient boosted trees", style="List Bullet")
    document.add_paragraph("Support vector machines", style="List Bullet")
    document.add_paragraph(
        "Regularisation reduces overfitting by penalising large coefficients."
    )

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Model"
    table.cell(0, 1).text = "Accuracy"
    table.cell(1, 0).text = "Random forest"
    table.cell(1, 1).text = "0.91"

    document.save(str(path))
    return path


def make_pptx() -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    path = FIXTURES / "sample.pptx"
    presentation = Presentation()

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Retrieval Augmented Generation"
    body = slide.placeholders[1].text_frame
    body.text = "Retrieve relevant context"
    body.add_paragraph().text = "Ground the model in real sources"
    slide.notes_slide.notes_text_frame.text = "Remember to mention citation accuracy."

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "Evaluation"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Measure retrieval recall"
    body2.add_paragraph().text = "Track answer faithfulness"

    slide3 = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide3.shapes.title.text = "Architecture"
    box = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    box.text_frame.text = "Embeddings are stored in PostgreSQL with pgvector"

    presentation.save(str(path))
    return path


def make_legacy_office() -> Path | None:
    """Real .doc and .ppt files, produced by LibreOffice itself.

    Hand-written OLE2 stubs would not exercise the converter, so the fixtures
    are genuine legacy documents converted from the modern ones. Skipped when
    LibreOffice is unavailable.
    """
    import subprocess
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
    import capabilities

    soffice = capabilities.libreoffice_path()
    if not soffice:
        print("  skipping .doc/.ppt fixtures: LibreOffice not found", file=sys.stderr)
        return None

    for source, target in [("sample.docx", "doc"), ("sample.pptx", "ppt")]:
        source_path = FIXTURES / source
        if not source_path.exists():
            continue
        subprocess.run(
            [
                soffice,
                "--headless",
                "--norestore",
                "--convert-to",
                target,
                "--outdir",
                str(FIXTURES),
                str(source_path),
            ],
            capture_output=True,
            timeout=180,
            check=False,
        )

    return FIXTURES / "sample.doc"


def make_unsupported() -> Path:
    path = FIXTURES / "script.sh"
    path.write_text("#!/bin/sh\necho hello\n", encoding="utf-8")
    return path


def make_corrupted_docx() -> Path:
    """Valid ZIP magic bytes, but not a real DOCX package."""
    path = FIXTURES / "corrupted.docx"
    path.write_bytes(b"PK\x03\x04" + b"\x00" * 64 + b"not really a docx package")
    return path


def make_corrupted_pdf() -> Path:
    path = FIXTURES / "corrupted.pdf"
    path.write_bytes(b"%PDF-1.7\n" + b"\xde\xad\xbe\xef" * 40)
    return path


def make_empty_txt() -> Path:
    path = FIXTURES / "empty.txt"
    path.write_bytes(b"")
    return path


def make_whitespace_txt() -> Path:
    path = FIXTURES / "whitespace.txt"
    path.write_text("   \n\n\t\n   \n", encoding="utf-8")
    return path


def make_fake_pdf() -> Path:
    """A .pdf extension on content that is really a shell script."""
    path = FIXTURES / "fake.pdf"
    path.write_bytes(b"#!/bin/sh\nrm -rf /\n")
    return path


BUILDERS = [
    make_txt,
    make_md,
    make_pdf_text,
    make_pdf_scanned,
    make_pdf_mixed,
    make_docx,
    make_pptx,
    make_legacy_office,
    make_unsupported,
    make_corrupted_docx,
    make_corrupted_pdf,
    make_empty_txt,
    make_whitespace_txt,
    make_fake_pdf,
]


# Paragraph used to fill the multi-page scalability fixtures. Roughly two
# 1000-character chunks per page once the splitter has run.
_FILLER = (
    "Retrieval augmented generation grounds a language model in retrieved "
    "context so that answers can cite real sources rather than relying only "
    "on parametric memory. Chunking, embedding and vector search each shape "
    "the quality of the retrieved context in different ways. "
)


def make_paged_pdf(pages: int, filename: str) -> Path:
    """A text PDF of `pages` pages, rebuilt only when the page count differs.

    Used by the memory tests, which need documents long enough that any
    per-page accumulation would be obvious. These are generated rather than
    committed: at 1000 pages the file is far too large to keep in the repo.
    """
    import pymupdf

    path = FIXTURES / filename
    if path.exists():
        existing = pymupdf.open(str(path))
        count = existing.page_count
        existing.close()
        if count == pages:
            return path

    FIXTURES.mkdir(parents=True, exist_ok=True)
    document = pymupdf.open()
    for number in range(pages):
        page = document.new_page()
        page.insert_textbox(
            pymupdf.Rect(40, 40, 560, 780),
            f"Page {number + 1}\n" + _FILLER * 6,
            fontsize=9,
        )
    document.save(str(path))
    document.close()
    return path


def build_all() -> list[Path]:
    FIXTURES.mkdir(parents=True, exist_ok=True)
    created = []
    for builder in BUILDERS:
        try:
            created.append(builder())
        except Exception as exc:  # noqa: BLE001
            print(f"  could not build {builder.__name__}: {exc}", file=sys.stderr)
    return created


if __name__ == "__main__":
    for item in build_all():
        print(f"  {item.name}  ({item.stat().st_size} bytes)")
